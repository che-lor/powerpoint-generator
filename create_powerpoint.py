#!/usr/bin/env python3

from pptx import Presentation
import itertools

###################-FILL-OUT-####################
file1 = "name_of_file"
file2 = "name_of_file"
file3 = "name_of_file"
powerpoint_name = "name_of_powerpoint"
#################################################
_file1 = file1 + ".txt"
_file2 = file2 + ".txt"
_file3 = file3 + ".txt"
_powerpoint_name = powerpoint_name + ".pptx"
slide_master_template = "slide_master.pptx"
save_message = "Powerpoint was created!"
error_message = "Looks like there was an error :("

with open(_file1) as f1, open(_file2) as f2, open(_file3) as f3:
	file1_lines = f1.read().splitlines()
	file2_lines = f2.read().splitlines()	
	file3_lines = f3.read().splitlines()
	lines = list(file1_lines + file2_lines + file3_lines)
	f1.close()
	f2.close()
	f3.close()

#Opening presentation using the "slide_master.pptx" file
prs = Presentation(slide_master_template)

#Create a slide for each line in the file
cycle = itertools.cycle(lines)
for eachline in lines:
	next_line = next(cycle)
	bullet_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	body_shape = shapes.placeholders[1]
	tf = body_shape.text_frame
	tf.text = next_line

save = prs.save(_powerpoint_name)
save
if [ save ]:
	print(save_message)
else:
	print(error_message)
