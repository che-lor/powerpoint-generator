#!/usr/bin/env python3

from pptx import Presentation
import itertools

###################-FILL-OUT-####################
file1 = "name_of_file"
file2 = "name_of_file"
file3 = "name_of_file"
powerpoint_name = "name_of_powerpoint"
#################################################

files = [powerpoint_name, file1, file2, file3]
ext = [".pptx", ".txt"]

#Store fullname of files
file_names = []
for f in files:
	if f == powerpoint_name:
		file_names.append(f + ext[0])
	else:
		file_names.append(f + ext[1])
# print(file_names)

lines = []
title_cycle = itertools.cycle(files[1:])
for fn in file_names[1:]:
	file = open(fn, 'r')
	file_lines = file.read().splitlines()
	lines.append(next(title_cycle))
	lines.append("")
	lines.extend(file_lines)
	if ((int(len(file_lines))) % 2) != 0:
		lines.append("")
		print("Added extra line")
	file.close()
# print(lines)

prs = Presentation("slide_master.pptx")

LPS = 2
total_slides = int(len(lines)/LPS)
line_cycle = itertools.cycle(lines)
loop = 0
while loop < total_slides:
	pair = next(line_cycle) + '\n' + next(line_cycle) # add another cycle for more LPS
	bullet_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	body_shape = shapes.placeholders[1]
	tf = body_shape.text_frame
	tf.text = pair
	loop += 1
	if loop > total_slides:
		break

#Save file
save = prs.save(file_names[0])
if [ save ]:
	print(file_names[0] + "was created!")
