#!/usr/bin/env python3

from pptx import Presentation
import itertools

###################-FILL-OUT-####################
file1 = "name_of_file"
file2 = "name_of_file"
file3 = "name_of_file"
powerpoint_name = "name_of_powerpoint"
#################################################
titles = [file1, file2, file3]
_file1 = file1 + ".txt"
_file2 = file2 + ".txt"
_file3 = file3 + ".txt"
_powerpoint_name = powerpoint_name + ".pptx"
slide_master_template = "slide_master.pptx"
save_message = _powerpoint_name + " was created!"

#Open & Read files to separate lists
with open(_file1) as f1, open(_file2) as f2, open(_file3) as f3:
	file1_lines = list(f1.read().splitlines())
	file2_lines = list(f2.read().splitlines())	
	file3_lines = list(f3.read().splitlines())
	#Close files
	f1.close()
	f2.close()
	f3.close()
	#Insert Titles in front of each list
	file1_lines.insert(0, file1)
	file2_lines.insert(0, file2)
	file3_lines.insert(0, file3)
	#Insert blank line after Titles
	file1_lines.insert(1, "")
	file2_lines.insert(1, "")
	file3_lines.insert(1, "")
	#Append blank lines if count is odd
	if ((int(len(file1_lines))) % 2) != 0:
		print(file1 + " line count is odd (" + str(len(file1_lines)) + ")...appended extra line.")
		file1_lines.append("")
	elif ((int(len(file2_lines))) % 2) != 0:
		print(file2 + " line count is odd (" + str(len(file2_lines)) + ")...appended extra line.")
		file2_lines.append("")
	elif ((int(len(file3_lines))) % 2) != 0:
		print(file3 + " line count is odd (" + str(len(file3_lines)) + ")...appended extra line.")
		file3_lines.append("")

#Open presentation using the "slide_master.pptx" file
prs = Presentation(slide_master_template)

#Create slides for each LPS
lines = file1_lines + file2_lines + file3_lines
LPS = 2
total_slides = int(len(lines)/LPS)
cycle = itertools.cycle(lines)
loop = 0
while loop < total_slides:
	pair = next(cycle) + '\n' + next(cycle) # add another cycle for more LPS
	bullet_slide_layout = prs.slide_layouts[1]
	slide = prs.slides.add_slide(bullet_slide_layout)
	shapes = slide.shapes
	body_shape = shapes.placeholders[1]
	tf = body_shape.text_frame
	tf.text = pair
	loop += 1
	if loop > total_slides:
		break

#Save file
save = prs.save(_powerpoint_name)
if [ save ]:
	print(save_message)
